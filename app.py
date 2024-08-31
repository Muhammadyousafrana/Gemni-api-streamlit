# import streamlit as st
# import google.generativeai as genai
# import os
# import re
# from pptx.util import Pt
# from pptx.dml.color import RGBColor
# from pptx import Presentation
# #from dotenv import load_dotenv
# #import json

# # Load environment variables
# #load_dotenv()
# #api_key = os.getenv('my_api_key')
# api_key=st.secrets["api_key"]
# # Configure the API key
# genai.configure(api_key=api_key)


# # Function to get response from Gemini API
# def get_gemini_response(prompt):
#     model = genai.GenerativeModel("gemini-1.5-pro")
#     response = model.generate_content([prompt])
#     return response.text


# # Function to clean text for PowerPoint slides
# def clean_text(text):
#     cleaned_text = re.sub(r'\s+', ' ', text).strip()
#     cleaned_text = re.sub(r'[*-]\s*|\d+\.\s*', '', cleaned_text)
#     cleaned_text = re.sub(r'\s*:\s*', ': ', cleaned_text)
#     cleaned_text = re.sub(r'\s*-\s*', ' - ', cleaned_text)
#     return cleaned_text


# # Function to split text into sentences
# def split_sentences(text):
#     sentences = re.split(r'(?<=\.)\s+', text)
#     sentences = [sentence.capitalize() for sentence in sentences]
#     return sentences


# # Function to replace and capitalize text between colons
# def replace_and_capitalize(text):
#     def replace_and_capitalize_colon(match):
#         return match.group(1) + match.group(2).capitalize() + match.group(3)

#     result = re.sub(r'(:\s*)(.*?)(\s*:[^:]|$)', replace_and_capitalize_colon, text)
#     return result


# # Function to refine content for PowerPoint slides
# def refine_final_content(content):
#     final_content = []
#     for i in content:
#         cleaned_text = clean_text(i)
#         sentences = split_sentences(cleaned_text)
#         final_content.append(sentences)
#     return final_content


# # Function to create PowerPoint slides
# def slide_maker(powerpoint, topic, sub_titles, final_content):
#     title_slide_layout = powerpoint.slide_layouts[0]
#     title_slide = powerpoint.slides.add_slide(title_slide_layout)
#     title = title_slide.shapes.title
#     title.text = topic
#     title.text_frame.paragraphs[0].font.size = Pt(32)
#     title.text_frame.paragraphs[0].font.bold = True
#     content = title_slide.placeholders[1]
#     content.text = "Created By AI Gemini Model"

#     for i in range(len(sub_titles)):
#         bulletLayout = powerpoint.slide_layouts[1]
#         secondSlide = powerpoint.slides.add_slide(bulletLayout)
#         myShapes = secondSlide.shapes
#         titleShape = myShapes.title
#         bodyShape = myShapes.placeholders[1]
#         titleShape.text = sub_titles[i]
#         titleShape.text_frame.paragraphs[0].font.size = Pt(24)
#         titleShape.text_frame.paragraphs[0].font.bold = True
#         tFrame = bodyShape.text_frame

#         for point in final_content[i]:
#             point = re.sub(r':[^:]+:', ':', point)
#             point = replace_and_capitalize(point)
#             p = tFrame.add_paragraph()
#             p.text = point
#             p.font.size = Pt(18)
#             p.space_after = Pt(10)

#     return powerpoint


# # Function to provide download button for PowerPoint
# def download_button(file_path, topic):
#     with open(file_path, "rb") as file:
#         ppt_content = file.read()

#     st.download_button(
#         label="Download PowerPoint",
#         data=ppt_content,
#         file_name=f"{topic}.pptx",
#         key="ppt_download_button"
#     )


# # Streamlit interface
# st.set_page_config(page_title="AI Content Generator")

# st.title("AI Content Generator")
# st.write("Select the tool you want to use:")

# option = st.selectbox(
#     "Choose a feature:",
#     ("Generate Blog Post", "Create Presentation", "Generate Product Description")
# )

# if option == "Generate Blog Post":
#     st.header("AI Blog Post Generator")
#     st.write("This tool uses the Gemini AI to generate blog posts based on your prompt.")

#     prompt = st.text_area("Enter your blog post idea or prompt:", height=200)

#     if st.button("Generate Blog Post"):
#         if prompt.strip():
#             st.write("Generating blog post, please wait...")
#             blog_post = get_gemini_response(prompt)
#             st.subheader("Generated Blog Post")
#             st.write(blog_post)
#         else:
#             st.write("Please enter a valid prompt to generate a blog post.")

# elif option == "Create Presentation":
#     st.header("Gemini Presentation Maker")

#     topic = st.text_input("Input Prompt:")
#     no_of_slide = st.text_input("Enter Number Of Slides:")

#     if st.button("Generate Presentation"):
#         if topic.strip() and no_of_slide.isdigit():
#             prompt = f"Generate {no_of_slide} sub-titles on the topic of {topic}"
#             response = get_gemini_response(prompt)

#             sub_topics = response.split("\n")
#             sub_titles = [sub_topic[3:].replace('"', "") for sub_topic in sub_topics]

#             content = []
#             for title in sub_titles:
#                 content_prompt = f"Generate content for {title} for presentation slides with 2 bullet points, each of 20 tokens"
#                 slide_content = get_gemini_response(content_prompt)
#                 content.append(slide_content)

#             final_content = refine_final_content(content)
#             powerpoint = Presentation()
#             powerpoint = slide_maker(powerpoint, topic, sub_titles, final_content)

#             file_path = f"{topic}.pptx"
#             powerpoint.save(file_path)
#             st.success("Presentation is ready!")
#             download_button(file_path, topic)
#         else:
#             st.write("Please provide a valid topic and number of slides.")

# elif option == "Generate Product Description":
#     st.header("Product Description Writer")
#     st.write("Enter product details below to generate a compelling product description:")

#     product_details = st.text_area("Product Details", placeholder="Enter product features, specifications, etc.")

#     if st.button("Generate Description"):
#         if product_details.strip():
#             model_name = 'gemini-1.5-flash'  # You can modify this based on the JSON file or make it dynamic
#             description = get_gemini_response(product_details)
#             st.subheader("Generated Product Description")
#             st.write(description)
#         else:
#             st.write("Please enter product details to generate a description.")


import streamlit as st
import google.generativeai as genai
import os
import re
from pptx.util import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor

# Load environment variables (if needed)
# from dotenv import load_dotenv
# load_dotenv()

# Set the API key for the Gemini API
api_key = st.secrets["api_key"]
genai.configure(api_key=api_key)

# Function to get response from Gemini API with error handling
def get_gemini_response(prompt):
    try:
        model = genai.GenerativeModel("gemini-1.5-pro")
        response = model.generate_content([prompt])
        return response.text
    except google.api_core.exceptions.InternalServerError:
        st.error("The Gemini API encountered an internal server error. Please try again later.")
        return None
    except Exception as e:
        st.error(f"An unexpected error occurred: {str(e)}")
        return None

# Function to clean text for PowerPoint slides
def clean_text(text):
    cleaned_text = re.sub(r'\s+', ' ', text).strip()
    cleaned_text = re.sub(r'[*-]\s*|\d+\.\s*', '', cleaned_text)
    cleaned_text = re.sub(r'\s*:\s*', ': ', cleaned_text)
    cleaned_text = re.sub(r'\s*-\s*', ' - ', cleaned_text)
    return cleaned_text

# Function to split text into sentences
def split_sentences(text):
    sentences = re.split(r'(?<=\.)\s+', text)
    sentences = [sentence.capitalize() for sentence in sentences]
    return sentences

# Function to replace and capitalize text between colons
def replace_and_capitalize(text):
    def replace_and_capitalize_colon(match):
        return match.group(1) + match.group(2).capitalize() + match.group(3)
    result = re.sub(r'(:\s*)(.*?)(\s*:[^:]|$)', replace_and_capitalize_colon, text)
    return result

# Function to refine content for PowerPoint slides
def refine_final_content(content):
    final_content = []
    for i in content:
        cleaned_text = clean_text(i)
        sentences = split_sentences(cleaned_text)
        final_content.append(sentences)
    return final_content

# Function to create PowerPoint slides
def slide_maker(powerpoint, topic, sub_titles, final_content):
    title_slide_layout = powerpoint.slide_layouts[0]
    title_slide = powerpoint.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = topic
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    content = title_slide.placeholders[1]
    content.text = "Created By AI Gemini Model"

    for i in range(len(sub_titles)):
        bulletLayout = powerpoint.slide_layouts[1]
        slide = powerpoint.slides.add_slide(bulletLayout)
        titleShape = slide.shapes.title
        bodyShape = slide.placeholders[1]
        titleShape.text = sub_titles[i]
        titleShape.text_frame.paragraphs[0].font.size = Pt(24)
        titleShape.text_frame.paragraphs[0].font.bold = True
        tFrame = bodyShape.text_frame

        for point in final_content[i]:
            point = re.sub(r':[^:]+:', ':', point)
            point = replace_and_capitalize(point)
            p = tFrame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.space_after = Pt(10)

    return powerpoint

# Function to provide download button for PowerPoint
def download_button(file_path, topic):
    with open(file_path, "rb") as file:
        ppt_content = file.read()

    st.download_button(
        label="Download PowerPoint",
        data=ppt_content,
        file_name=f"{topic}.pptx",
        key="ppt_download_button"
    )

# Streamlit interface
st.set_page_config(page_title="AI Content Generator")

st.title("AI Content Generator")
st.write("Select the tool you want to use:")

option = st.selectbox(
    "Choose a feature:",
    ("Generate Blog Post", "Create Presentation", "Generate Product Description")
)

if option == "Generate Blog Post":
    st.header("AI Blog Post Generator")
    st.write("This tool uses the Gemini AI to generate blog posts based on your prompt.")

    prompt = st.text_area("Enter your blog post idea or prompt:", height=200)

    if st.button("Generate Blog Post"):
        if prompt.strip():
            st.write("Generating blog post, please wait...")
            blog_post = get_gemini_response(prompt)
            if blog_post:
                st.subheader("Generated Blog Post")
                st.write(blog_post)
        else:
            st.write("Please enter a valid prompt to generate a blog post.")

elif option == "Create Presentation":
    st.header("Gemini Presentation Maker")

    topic = st.text_input("Input Prompt:")
    no_of_slide = st.text_input("Enter Number Of Slides:")

    if st.button("Generate Presentation"):
        if topic.strip() and no_of_slide.isdigit():
            prompt = f"Generate {no_of_slide} sub-titles on the topic of {topic}"
            response = get_gemini_response(prompt)

            if response:
                sub_topics = response.split("\n")
                sub_titles = [sub_topic[3:].replace('"', "") for sub_topic in sub_topics]

                content = []
                for title in sub_titles:
                    content_prompt = f"Generate content for {title} for presentation slides with 2 bullet points, each of 20 tokens"
                    slide_content = get_gemini_response(content_prompt)
                    if slide_content:
                        content.append(slide_content)

                if content:
                    final_content = refine_final_content(content)
                    powerpoint = Presentation()
                    powerpoint = slide_maker(powerpoint, topic, sub_titles, final_content)

                    file_path = f"{topic}.pptx"
                    powerpoint.save(file_path)
                    st.success("Presentation is ready!")
                    download_button(file_path, topic)
        else:
            st.write("Please provide a valid topic and number of slides.")

elif option == "Generate Product Description":
    st.header("Product Description Writer")
    st.write("Enter product details below to generate a compelling product description:")

    product_details = st.text_area("Product Details", placeholder="Enter product features, specifications, etc.")

    if st.button("Generate Description"):
        if product_details.strip():
            description = get_gemini_response(product_details)
            if description:
                st.subheader("Generated Product Description")
                st.write(description)
        else:
            st.write("Please enter product details to generate a description.")

